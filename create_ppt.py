from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─── COLOR PALETTE ────────────────────────────────────────────
BROWN      = RGBColor(0x8C, 0x52, 0x26)
BROWN_DARK = RGBColor(0x5A, 0x32, 0x14)
YELLOW     = RGBColor(0xF5, 0xA6, 0x23)
RED        = RGBColor(0xE6, 0x39, 0x46)
BLACK      = RGBColor(0x11, 0x11, 0x11)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CREAM      = RGBColor(0xF7, 0xF5, 0xF0)
LIGHT_GRAY = RGBColor(0xE8, 0xE5, 0xDE)
BROWN_TEXT = RGBColor(0x5C, 0x54, 0x4D)

# ─── HELPERS ──────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return txBox

def add_para(tf, text, size, color, bold=False, level=0, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Calibri"
    return p

def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def side_accent(slide, color=BROWN):
    add_rect(slide, 0, 0, 0.35, 7.5, color)

def top_bar(slide, color=BROWN):
    add_rect(slide, 0, 0, 13.33, 0.12, color)
    add_rect(slide, 0, 7.38, 13.33, 0.12, color)

def slide_number(slide, num, total=9):
    add_text(slide, f"{num} / {total}", 11.8, 7.1, 1.5, 0.4, 9, BROWN_TEXT, align=PP_ALIGN.RIGHT)

def section_tag(slide, label, left=0.6, top=0.25):
    tag = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(2.2), Inches(0.32))
    tag.fill.solid()
    tag.fill.fore_color.rgb = YELLOW
    tag.line.fill.background()
    tf = tag.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label.upper()
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.name = "Calibri"
    run.font.color.rgb = BLACK

def divider_line(slide, left, top, width, color=LIGHT_GRAY):
    ln = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.025))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()

def bullet_box(slide, items, left, top, width, height, dot_color=BROWN, font_size=14):
    """Items: list of (text, is_sub) tuples"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for text, is_sub in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4 if not is_sub else 2)
        p.level = 1 if is_sub else 0
        run = p.add_run()
        prefix = "    ◦  " if is_sub else "▸  "
        run.text = prefix + text
        run.font.size = Pt(font_size - 1 if is_sub else font_size)
        run.font.color.rgb = BROWN_TEXT if is_sub else BLACK
        run.font.name = "Calibri"
        run.font.bold = not is_sub

def info_card(slide, title, value, left, top, width=2.8, height=1.4, accent=BROWN):
    card = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.75)
    # top accent strip
    strip = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.07))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    add_text(slide, value, left+0.15, top+0.15, width-0.3, 0.6, 22, accent, bold=True)
    add_text(slide, title, left+0.15, top+0.75, width-0.3, 0.5, 11, BROWN_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
slide = blank_slide()
# Full dark background
add_rect(slide, 0, 0, 13.33, 7.5, BLACK)
# Warm accent overlay shapes
add_rect(slide, 0, 0, 6.0, 7.5, BROWN_DARK)
# Yellow accent bar
add_rect(slide, 5.7, 0, 0.2, 7.5, YELLOW)
# Logo area circle
circ = slide.shapes.add_shape(9, Inches(1.2), Inches(1.6), Inches(2.0), Inches(2.0))  # oval
circ.fill.solid()
circ.fill.fore_color.rgb = YELLOW
circ.line.fill.background()
add_text(slide, "GC", 1.45, 1.9, 1.5, 1.0, 36, BLACK, bold=True, align=PP_ALIGN.CENTER)
# Main title
add_text(slide, "GlobalCart", 0.5, 3.8, 5.0, 1.1, 46, WHITE, bold=True)
add_text(slide, "Digital Products Marketplace", 0.5, 4.75, 5.1, 0.7, 18, YELLOW)
add_text(slide, "DBMS Mini Project  —  2025–26", 0.5, 5.4, 5.1, 0.5, 12, RGBColor(0xCC,0xBB,0xAA))
# Right side — team
add_text(slide, "Team Members", 6.5, 1.5, 6.5, 0.5, 13, YELLOW, bold=True)
divider_line(slide, 6.5, 2.1, 6.3, BROWN)

members = [
    ("Aaditya Singh", "1TD24AI002"),
    ("Harsha Vardhan", "1BY25AI411"),
    ("Vishal Kumar", "1TD24AI188"),
    ("Chalapathi", "1BY25AI411"),
]
for i, (name, usn) in enumerate(members):
    y = 2.3 + i * 1.05
    add_rect(slide, 6.4, y, 6.4, 0.82, RGBColor(0x22, 0x1A, 0x14))
    add_rect(slide, 6.4, y, 0.08, 0.82, YELLOW)
    add_text(slide, name, 6.65, y + 0.05, 5.5, 0.4, 15, WHITE, bold=True)
    add_text(slide, usn, 6.65, y + 0.43, 5.5, 0.3, 11, YELLOW)

add_text(slide, "Database Management Systems Laboratory", 6.4, 6.8, 6.8, 0.4, 10, RGBColor(0x99,0x88,0x77), align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCTION
# ══════════════════════════════════════════════════════════════
slide = blank_slide()
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
side_accent(slide, BROWN)
top_bar(slide)
section_tag(slide, "01 · Introduction")
add_text(slide, "Introduction & Problem Statement", 0.6, 0.7, 11.5, 0.8, 30, BLACK, bold=True)
divider_line(slide, 0.6, 1.6, 12.5, BROWN)
slide_number(slide, 2)

points = [
    ("What is GlobalCart?", False),
    ("A modern e-commerce platform for digital products — e-books, courses, software tools.", True),
    ("Objective", False),
    ("Build a robust, scalable database-driven web application handling authentication, product catalog, cart and order processing end-to-end.", True),
    ("Why We Built This?", False),
    ("Traditional e-commerce is designed for physical goods. GlobalCart provides instant delivery relying entirely on a structured relational database.", True),
    ("Core Challenge", False),
    ("Designing normalized schemas that maintain data integrity across users, products, carts, and transactions.", True),
]
bullet_box(slide, points, 0.65, 1.75, 12.4, 5.4, font_size=14)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — TECH STACK
# ══════════════════════════════════════════════════════════════
slide = blank_slide()
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
side_accent(slide, BROWN)
top_bar(slide)
section_tag(slide, "02 · Architecture")
add_text(slide, "Technology Stack & System Architecture", 0.6, 0.7, 11.5, 0.8, 30, BLACK, bold=True)
divider_line(slide, 0.6, 1.6, 12.5, BROWN)
slide_number(slide, 3)

# 3 info cards row 1
info_card(slide, "Frontend", "HTML5 + CSS + JS", 0.65, 1.8, 3.5, 1.3, BROWN)
info_card(slide, "Backend / API", "Python + Flask", 4.55, 1.8, 3.5, 1.3, RED)
info_card(slide, "Database", "MySQL (RDBMS)", 8.45, 1.8, 3.5, 1.3, BROWN_DARK)

info_card(slide, "Architecture", "SPA + REST API", 0.65, 3.4, 3.5, 1.3, YELLOW)
info_card(slide, "Security", "Hashed Passwords + Sessions", 4.55, 3.4, 3.5, 1.3, BROWN)
info_card(slide, "Connection", "MySQL Connection Pooling", 8.45, 3.4, 3.5, 1.3, RED)

add_text(slide, "Architecture Flow:", 0.65, 4.95, 12.0, 0.35, 12, BLACK, bold=True)
add_text(slide, "Browser (SPA)  →  Flask REST API (Python)  →  MySQL Connection Pool  →  MySQL Database", 0.65, 5.35, 12.0, 0.5, 13, BROWN_TEXT)
divider_line(slide, 0.65, 5.95, 12.3, LIGHT_GRAY)
add_text(slide, "All API endpoints are protected. The frontend never directly touches the database — all operations go through Flask.", 0.65, 6.1, 12.0, 0.6, 11, BROWN_TEXT, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — E-R DESIGN
# ══════════════════════════════════════════════════════════════
slide = blank_slide()
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
side_accent(slide, RED)
top_bar(slide, RED)
section_tag(slide, "03 · ER Design")
add_text(slide, "Entity-Relationship (E-R) Design", 0.6, 0.7, 11.5, 0.8, 30, BLACK, bold=True)
divider_line(slide, 0.6, 1.6, 12.5, RED)
slide_number(slide, 4)

entities = [
    ("Users", "Stores customer & admin credentials. Attributes: id, username, email (UNIQUE), password_hash, role.", BROWN),
    ("Products", "Digital item metadata. Attributes: id, title, description, price, category, image_url, file_path.", RED),
    ("Cart", "Temporary storage before purchase. Attributes: id, user_id (FK), product_id (FK), quantity.", BROWN_DARK),
    ("Orders", "Confirmed purchases. Attributes: id, user_id (FK), total_amount, payment_method, status.", YELLOW),
    ("Order_Items", "Resolves Many-to-Many between Orders and Products. Stores price_at_time of purchase.", BROWN),
]

for i, (name, desc, color) in enumerate(entities):
    row = i // 2
    col = i % 2
    if i == 4:
        lft = 0.65 + 4 * 2.3
        tp = 1.85 + row * 2.1
    else:
        lft = 0.65 + col * 6.15
        tp = 1.85 + row * 2.1
    card = slide.shapes.add_shape(1, Inches(lft), Inches(tp), Inches(5.8), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(1)
    strip = slide.shapes.add_shape(1, Inches(lft), Inches(tp), Inches(5.8), Inches(0.1))
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.fill.background()
    add_text(slide, name, lft + 0.15, tp + 0.15, 5.4, 0.45, 16, color, bold=True)
    add_text(slide, desc, lft + 0.15, tp + 0.65, 5.5, 1.0, 10, BROWN_TEXT, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — SCHEMA
# ══════════════════════════════════════════════════════════════
slide = blank_slide()
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
side_accent(slide, YELLOW)
top_bar(slide, YELLOW)
section_tag(slide, "04 · Schema")
add_text(slide, "Relational Schema & Database Tables", 0.6, 0.7, 11.5, 0.8, 30, BLACK, bold=True)
divider_line(slide, 0.6, 1.6, 12.5, YELLOW)
slide_number(slide, 5)

tables = [
    ("users", "id (PK)  ·  username  ·  email (UNIQUE)  ·  password_hash  ·  role", BROWN),
    ("products", "id (PK)  ·  title  ·  description  ·  price  ·  category  ·  image_url  ·  file_path", RED),
    ("cart", "id (PK)  ·  user_id (FK→users)  ·  product_id (FK→products)  ·  quantity", BROWN_DARK),
    ("orders", "id (PK)  ·  user_id (FK→users)  ·  total_amount  ·  payment_method  ·  status", YELLOW),
    ("order_items", "id (PK)  ·  order_id (FK→orders)  ·  product_id (FK→products)  ·  price_at_time", BROWN),
]

for i, (tname, cols, color) in enumerate(tables):
    y = 1.85 + i * 0.95
    add_rect(slide, 0.65, y, 12.3, 0.75, WHITE)
    acc = slide.shapes.add_shape(1, Inches(0.65), Inches(y), Inches(0.18), Inches(0.75))
    acc.fill.solid()
    acc.fill.fore_color.rgb = color
    acc.line.fill.background()
    add_text(slide, tname, 1.0, y + 0.06, 2.8, 0.4, 13, color, bold=True)
    add_text(slide, cols, 3.8, y + 0.1, 9.0, 0.55, 10, BROWN_TEXT)

add_text(slide, "Key Constraints Used", 0.65, 6.75, 5.0, 0.4, 11, BLACK, bold=True)
add_text(slide, "PRIMARY KEY  ·  FOREIGN KEY (Referential Integrity)  ·  UNIQUE  ·  ON DELETE CASCADE  ·  NOT NULL", 0.65, 7.05, 12.3, 0.3, 10, BROWN_TEXT)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — KEY FEATURES
# ══════════════════════════════════════════════════════════════
slide = blank_slide()
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
side_accent(slide, BROWN)
top_bar(slide)
section_tag(slide, "05 · Features")
add_text(slide, "Key Features Demonstrated", 0.6, 0.7, 11.5, 0.8, 30, BLACK, bold=True)
divider_line(slide, 0.6, 1.6, 12.5, BROWN)
slide_number(slide, 6)

features = [
    ("🔐", "Secure Authentication", "User registration & login using hashed passwords. Role-based access (Admin / Customer).", BROWN),
    ("🛍️", "Dynamic Product Catalog", "Products fetched live from MySQL. Supports category filtering, search, and real-time listing.", RED),
    ("🛒", "Cart Management", "Real-time DB updates when adding/removing items. Persistent cart across sessions.", BROWN_DARK),
    ("💳", "Checkout & Orders", "Transactional inserts generating Order IDs and linking Order Items with price_at_time.", YELLOW),
    ("📊", "Admin Dashboard", "Aggregation queries (SUM, COUNT) to display total sales, user count, and recent orders.", BROWN),
    ("📦", "Instant Download", "On successful order, download links are generated and stored in the database.", RED),
]

for i, (icon, title, desc, color) in enumerate(features):
    row = i // 2
    col = i % 2
    lft = 0.6 + col * 6.35
    tp = 1.9 + row * 1.7
    card = slide.shapes.add_shape(1, Inches(lft), Inches(tp), Inches(6.0), Inches(1.45))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LIGHT_GRAY
    card.line.width = Pt(0.5)
    strip = slide.shapes.add_shape(1, Inches(lft), Inches(tp), Inches(6.0), Inches(0.07))
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.fill.background()
    add_text(slide, icon + "  " + title, lft + 0.15, tp + 0.12, 5.6, 0.45, 14, color, bold=True)
    add_text(slide, desc, lft + 0.15, tp + 0.65, 5.6, 0.7, 10, BROWN_TEXT, wrap=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — IMPLEMENTATION HIGHLIGHTS
# ══════════════════════════════════════════════════════════════
slide = blank_slide()
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
side_accent(slide, RED)
top_bar(slide, RED)
section_tag(slide, "06 · Implementation")
add_text(slide, "Implementation Highlights", 0.6, 0.7, 11.5, 0.8, 30, BLACK, bold=True)
divider_line(slide, 0.6, 1.6, 12.5, RED)
slide_number(slide, 7)

points = [
    ("MySQL Connection Pooling", False),
    ("Used mysql.connector.pooling to manage 10 simultaneous DB connections efficiently — prevents crashes under concurrent load.", True),
    ("Secure Password Storage", False),
    ("All passwords are hashed using Werkzeug's generate_password_hash before storing in the database. Plain-text passwords are never saved.", True),
    ("RESTful API Design", False),
    ("All features (auth, cart, products, orders) are separate Flask Blueprints exposing clean REST endpoints (GET / POST / DELETE).", True),
    ("Session-Based Authentication", False),
    ("Login sessions are server-side managed. API endpoints verify identity before any database read or write operation.", True),
    ("Transactional Integrity", False),
    ("Checkout uses a single DB transaction — if any insert fails (order or order_items), the entire transaction rolls back to maintain consistency.", True),
]
bullet_box(slide, points, 0.65, 1.75, 12.4, 5.3, font_size=13)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — CONCLUSION
# ══════════════════════════════════════════════════════════════
slide = blank_slide()
add_rect(slide, 0, 0, 13.33, 7.5, CREAM)
side_accent(slide, BROWN)
top_bar(slide)
section_tag(slide, "07 · Conclusion")
add_text(slide, "Conclusion & Future Scope", 0.6, 0.7, 11.5, 0.8, 30, BLACK, bold=True)
divider_line(slide, 0.6, 1.6, 12.5, BROWN)
slide_number(slide, 8)

# left: conclusion
add_text(slide, "What We Achieved", 0.65, 1.75, 5.8, 0.45, 14, BROWN, bold=True)
divider_line(slide, 0.65, 2.25, 5.8, BROWN)
c_points = [
    ("Successfully integrated MySQL relational DB with a modern web interface.", False),
    ("Implemented all core DBMS concepts: normalization, referential integrity, transactions.", False),
    ("Built a fully functional multi-user marketplace with Admin and Customer roles.", False),
    ("Applied connection pooling, secure auth, and session management in production-quality code.", False),
]
bullet_box(slide, c_points, 0.65, 2.4, 5.8, 3.5, font_size=12)

# divider
add_rect(slide, 6.7, 1.75, 0.04, 5.0, LIGHT_GRAY)

# right: future
add_text(slide, "Future Enhancements", 7.0, 1.75, 5.8, 0.45, 14, RED, bold=True)
divider_line(slide, 7.0, 2.25, 5.8, RED)
f_points = [
    ("Add DB Triggers for automated inventory tracking on purchase.", False),
    ("Implement FULLTEXT indexing for fast product search.", False),
    ("Integrate real payment gateway (Razorpay / Stripe).", False),
    ("Add product review & rating system with aggregate queries.", False),
    ("Cloud deployment with cloud-hosted MySQL for 24/7 availability.", False),
]
bullet_box(slide, f_points, 7.0, 2.4, 5.9, 3.5, font_size=12)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — THANK YOU
# ══════════════════════════════════════════════════════════════
slide = blank_slide()
add_rect(slide, 0, 0, 13.33, 7.5, BLACK)
add_rect(slide, 0, 0, 8.0, 7.5, BROWN_DARK)
add_rect(slide, 7.7, 0, 0.25, 7.5, YELLOW)

# Big GC
circ = slide.shapes.add_shape(9, Inches(1.5), Inches(1.5), Inches(2.4), Inches(2.4))
circ.fill.solid()
circ.fill.fore_color.rgb = YELLOW
circ.line.fill.background()
add_text(slide, "GC", 1.75, 1.85, 1.9, 1.4, 52, BLACK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, "Thank You!", 0.6, 4.1, 7.0, 1.2, 52, WHITE, bold=True)
add_text(slide, "We are now ready for the live demonstration.", 0.6, 5.25, 7.0, 0.55, 16, YELLOW)
add_text(slide, "Any Questions?", 0.6, 5.9, 7.0, 0.5, 14, RGBColor(0xCC, 0xBB, 0xAA), italic=True)

add_text(slide, "Team GlobalCart", 8.6, 2.5, 4.2, 0.5, 13, YELLOW, bold=True)
divider_line(slide, 8.6, 3.05, 4.0, BROWN)
for i, (name, usn) in enumerate(members):
    y = 3.25 + i * 0.82
    add_text(slide, name, 8.6, y, 4.0, 0.4, 13, WHITE, bold=True)
    add_text(slide, usn, 8.6, y + 0.38, 4.0, 0.3, 10, YELLOW)

output = r"c:\Users\AADI\Desktop\GlobalCart_DBMS_Presentation.pptx"
prs.save(output)
print(f"Saved: {output}")
